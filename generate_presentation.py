#!/usr/bin/env python3
"""Create a presentation from a template using data from a Word document."""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


FONT_NAME_REGULAR = "Montserrat Medium"
FONT_NAME_BOLD = "Montserrat Bold"


@dataclass
class QuestionItem:
    number: int | None
    theme: str
    question: str
    question_runs: list[dict[str, bool | str]] | None = None
    answer: str = ""
    comment: str = ""
    source: str = ""


FIELD_PATTERNS = {
    "theme": re.compile(r"^\s*(?:\d+\.\s*)?Тематика\s*:\s*(.*)$", re.IGNORECASE),
    "question": re.compile(r"^\s*Вопрос\s*:\s*(.*)$", re.IGNORECASE),
    "answer": re.compile(r"^\s*Ответ\s*:\s*(.*)$", re.IGNORECASE),
    "comment": re.compile(r"^\s*Комментарий\s*:\s*(.*)$", re.IGNORECASE),
    "source": re.compile(r"^\s*Источник\s*:\s*(.*)$", re.IGNORECASE),
}
NUMBER_PATTERN = re.compile(r"^\s*(\d+)\.\s*")


def extract_runs_slice(paragraph, start: int, end: int | None = None) -> list[dict[str, bool | str]]:
    chunks: list[dict[str, bool | str]] = []
    cursor = 0
    limit = len(paragraph.text) if end is None else end
    for run in paragraph.runs:
        run_text = run.text or ""
        run_len = len(run_text)
        run_start = cursor
        run_end = cursor + run_len
        cursor = run_end
        if run_len == 0 or run_end <= start or run_start >= limit:
            continue

        local_start = max(0, start - run_start)
        local_end = min(run_len, limit - run_start)
        text = run_text[local_start:local_end]
        if text:
            chunks.append(
                {
                    "text": text,
                    "bold": bool(run.bold),
                    "italic": bool(run.italic),
                    "underline": bool(run.underline),
                }
            )
    return chunks


def append_runs(target: list[dict[str, bool | str]], new_runs: list[dict[str, bool | str]], with_space: bool = False) -> list[dict[str, bool | str]]:
    if with_space and target:
        target.append({"text": " ", "bold": False, "italic": False, "underline": False})
    target.extend(new_runs)
    return target


def runs_to_plain_text(runs: list[dict[str, bool | str]] | None) -> str:
    if not runs:
        return ""
    return normalize_spaces("".join(str(item["text"]) for item in runs))


def is_field_line(line: str) -> bool:
    return any(pattern.match(line) for pattern in FIELD_PATTERNS.values())


def normalize_spaces(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def parse_questions_from_docx(docx_path: Path) -> list[QuestionItem]:
    """Parse questions from a .docx with blocks containing Тематика/Вопрос/... fields."""
    doc = Document(str(docx_path))

    items: list[QuestionItem] = []
    current: dict[str, str | int | None | list[dict[str, bool | str]]] | None = None
    current_field: str | None = None

    def flush_current() -> None:
        nonlocal current, current_field
        if not current:
            return
        theme = normalize_spaces(str(current.get("theme", "")))
        question_runs = current.get("question_runs")
        question = runs_to_plain_text(question_runs if isinstance(question_runs, list) else None)
        if not question:
            question = normalize_spaces(str(current.get("question", "")))
        if theme and question:
            items.append(
                QuestionItem(
                    number=current.get("number"),
                    theme=theme,
                    question=question,
                    question_runs=question_runs if isinstance(question_runs, list) else None,
                    answer=normalize_spaces(str(current.get("answer", ""))),
                    comment=normalize_spaces(str(current.get("comment", ""))),
                    source=normalize_spaces(str(current.get("source", ""))),
                )
            )
        current = None
        current_field = None

    for paragraph in doc.paragraphs:
        line = paragraph.text.strip()
        if not line:
            continue

        number_match = NUMBER_PATTERN.match(line)
        if number_match and FIELD_PATTERNS["theme"].match(line):
            flush_current()
            current = {"number": int(number_match.group(1))}
        elif number_match and not is_field_line(line):
            flush_current()
            current = None
            continue

        if current is None:
            current = {"number": None}

        matched_field = None
        for field_name, pattern in FIELD_PATTERNS.items():
            match = pattern.match(line)
            if match:
                matched_field = field_name
                value = match.group(1).strip()
                if value:
                    current[field_name] = value
                elif field_name not in current:
                    current[field_name] = ""
                current_field = field_name

                if field_name == "question":
                    field_start = match.start(1)
                    question_runs = extract_runs_slice(paragraph, field_start)
                    current["question_runs"] = question_runs
                break

        if matched_field:
            continue

        if current_field:
            existing = str(current.get(current_field, "")).strip()
            current[current_field] = f"{existing} {line}".strip()
            if current_field == "question":
                current_question_runs = current.get("question_runs")
                if isinstance(current_question_runs, list):
                    append_runs(current_question_runs, extract_runs_slice(paragraph, 0), with_space=True)

    flush_current()
    return items


def parse_round_without_theme_from_docx(
    docx_path: Path,
    round_title: str,
    max_questions: int,
) -> list[QuestionItem]:
    """Parse a numbered round where question text starts on lines like `1. ...` without `Тематика:`."""
    doc = Document(str(docx_path))
    lines = [p.text.strip() for p in doc.paragraphs]

    start_index = next((i for i, line in enumerate(lines) if normalize_spaces(line).lower() == round_title.lower()), None)
    if start_index is None:
        return []

    items: list[QuestionItem] = []
    current: dict[str, str | int | None] | None = None
    current_field: str | None = None

    def flush_current() -> None:
        nonlocal current, current_field
        if not current:
            return
        question = normalize_spaces(str(current.get("question", "")))
        if question:
            items.append(
                QuestionItem(
                    number=current.get("number"),
                    theme="",
                    question=question,
                    answer=normalize_spaces(str(current.get("answer", ""))),
                    comment=normalize_spaces(str(current.get("comment", ""))),
                    source=normalize_spaces(str(current.get("source", ""))),
                )
            )
        current = None
        current_field = None

    for raw_line in lines[start_index + 1 :]:
        line = raw_line.strip()
        if not line:
            continue

        if len(items) >= max_questions:
            break

        number_match = NUMBER_PATTERN.match(line)
        if number_match and not is_field_line(line):
            flush_current()
            current = {
                "number": int(number_match.group(1)),
                "question": normalize_spaces(NUMBER_PATTERN.sub("", line, count=1)),
            }
            current_field = "question"
            continue

        if current is None:
            continue

        matched_field = None
        for field_name, pattern in FIELD_PATTERNS.items():
            match = pattern.match(line)
            if match:
                matched_field = field_name
                if field_name == "theme":
                    # Для этого раунда тематику пропускаем.
                    current_field = None
                    break

                value = match.group(1).strip()
                if value:
                    current[field_name] = value
                elif field_name not in current:
                    current[field_name] = ""
                current_field = field_name
                break

        if matched_field:
            continue

        if current_field:
            existing = str(current.get(current_field, "")).strip()
            current[current_field] = f"{existing} {line}".strip()

    flush_current()
    return items[:max_questions]


def parse_round_with_section_themes_from_docx(
    docx_path: Path,
    round_title: str,
    max_questions: int,
) -> list[QuestionItem]:
    """Parse a numbered round where themes are declared as `Тематика N: ...` for groups of questions."""
    doc = Document(str(docx_path))
    lines = [p.text.strip() for p in doc.paragraphs]

    start_index = next((i for i, line in enumerate(lines) if normalize_spaces(line).lower() == round_title.lower()), None)
    if start_index is None:
        return []

    theme_section_pattern = re.compile(r"^\s*Тематика\s*\d+\s*:\s*(.*)$", re.IGNORECASE)

    items: list[QuestionItem] = []
    current_theme = ""
    current: dict[str, str | int | None] | None = None
    current_field: str | None = None

    def flush_current() -> None:
        nonlocal current, current_field
        if not current:
            return

        question = normalize_spaces(str(current.get("question", "")))
        if question:
            items.append(
                QuestionItem(
                    number=current.get("number"),
                    theme=normalize_spaces(str(current.get("theme", ""))),
                    question=question,
                    answer=normalize_spaces(str(current.get("answer", ""))),
                    comment=normalize_spaces(str(current.get("comment", ""))),
                    source=normalize_spaces(str(current.get("source", ""))),
                )
            )

        current = None
        current_field = None

    for raw_line in lines[start_index + 1 :]:
        line = raw_line.strip()
        if not line:
            continue

        if len(items) >= max_questions:
            break

        theme_match = theme_section_pattern.match(line)
        if theme_match:
            flush_current()
            current_theme = normalize_spaces(theme_match.group(1))
            continue

        number_match = NUMBER_PATTERN.match(line)
        if number_match and not is_field_line(line):
            flush_current()
            current = {
                "number": int(number_match.group(1)),
                "theme": current_theme,
                "question": normalize_spaces(NUMBER_PATTERN.sub("", line, count=1)),
            }
            current_field = "question"
            continue

        if current is None:
            continue

        matched_field = None
        for field_name, pattern in FIELD_PATTERNS.items():
            match = pattern.match(line)
            if match:
                matched_field = field_name
                if field_name == "theme":
                    current_field = None
                    break

                value = match.group(1).strip()
                if value:
                    current[field_name] = value
                elif field_name not in current:
                    current[field_name] = ""
                current_field = field_name
                break

        if matched_field:
            continue

        if current_field:
            existing = str(current.get(current_field, "")).strip()
            current[current_field] = f"{existing} {line}".strip()

    flush_current()
    return items[:max_questions]


def parse_round_with_constant_theme_from_docx(
    docx_path: Path,
    round_title: str,
    theme_label: str,
    max_questions: int,
) -> list[QuestionItem]:
    """Parse a numbered round with a fixed theme line like `Тема: 12`."""
    doc = Document(str(docx_path))
    lines = [p.text.strip() for p in doc.paragraphs]

    start_index = next((i for i, line in enumerate(lines) if normalize_spaces(line).lower() == round_title.lower()), None)
    if start_index is None:
        return []

    theme_value = ""
    theme_pattern = re.compile(rf"^\s*{re.escape(theme_label)}\s*:\s*(.*)$", re.IGNORECASE)
    for raw_line in lines[start_index + 1 :]:
        line = raw_line.strip()
        if not line:
            continue

        match = theme_pattern.match(line)
        if match:
            theme_value = normalize_spaces(match.group(1)).strip('"«»')
            break

    items: list[QuestionItem] = []
    current: dict[str, str | int | None] | None = None
    current_field: str | None = None

    def flush_current() -> None:
        nonlocal current, current_field
        if not current:
            return

        question = normalize_spaces(str(current.get("question", "")))
        if question:
            items.append(
                QuestionItem(
                    number=current.get("number"),
                    theme=theme_value,
                    question=question,
                    answer=normalize_spaces(str(current.get("answer", ""))),
                    comment=normalize_spaces(str(current.get("comment", ""))),
                    source=normalize_spaces(str(current.get("source", ""))),
                )
            )

        current = None
        current_field = None

    for raw_line in lines[start_index + 1 :]:
        line = raw_line.strip()
        if not line:
            continue

        if len(items) >= max_questions:
            break

        if theme_pattern.match(line):
            continue

        number_match = NUMBER_PATTERN.match(line)
        if number_match and not is_field_line(line):
            flush_current()
            current = {
                "number": int(number_match.group(1)),
                "question": normalize_spaces(NUMBER_PATTERN.sub("", line, count=1)),
            }
            current_field = "question"
            continue

        if current is None:
            continue

        matched_field = None
        for field_name, pattern in FIELD_PATTERNS.items():
            match = pattern.match(line)
            if match:
                matched_field = field_name
                if field_name == "theme":
                    current_field = None
                    break

                value = match.group(1).strip()
                if value:
                    current[field_name] = value
                elif field_name not in current:
                    current[field_name] = ""
                current_field = field_name
                break

        if matched_field:
            continue

        if current_field:
            existing = str(current.get(current_field, "")).strip()
            current[current_field] = f"{existing} {line}".strip()

    flush_current()
    return items[:max_questions]


def replace_placeholder(text: str, placeholder: str, value: str) -> str:
    return re.sub(re.escape(placeholder), value, text, flags=re.IGNORECASE)


def choose_question_font_size(question_text: str) -> int:
    length = len(question_text)
    if length <= 60:
        return 50
    if length <= 120:
        return 50
    if length <= 150:
        return 48
    if length <= 180:
        return 48
    if length <= 210:
        return 44
    if length <= 240:
        return 40
    return 36


def apply_font(run_obj, font_size: int, bold: bool = False, italic: bool = False, underline: bool = False) -> None:
    font = run_obj.font
    font.name = FONT_NAME_BOLD if bold else FONT_NAME_REGULAR
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.underline = underline


def fill_question_shape(shape, question_text: str, question_runs: list[dict[str, bool | str]] | None) -> None:
    font_size = choose_question_font_size(question_text)
    runs_data = question_runs or [{"text": question_text, "bold": False, "italic": False, "underline": False}]

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    for item in runs_data:
        r = p.add_run()
        r.text = str(item["text"])
        apply_font(
            r,
            font_size=font_size,
            bold=bool(item.get("bold", False)),
            italic=bool(item.get("italic", False)),
            underline=bool(item.get("underline", False)),
        )


def replace_in_text_frame(text_frame, replacements: dict[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            new_text = run.text
            for placeholder, value in replacements.items():
                new_text = replace_placeholder(new_text, placeholder, value)
            run.text = new_text



def replace_in_shape(shape, replacements: dict[str, str], question_runs: list[dict[str, bool | str]] | None = None) -> None:
    if getattr(shape, "has_text_frame", False):
        shape_text = normalize_spaces(shape.text_frame.text).lower()
        question_text = replacements.get("вопрос", "")
        if question_text and shape_text == "вопрос":
            fill_question_shape(shape, question_text, question_runs)
        else:
            replace_in_text_frame(shape.text_frame, replacements)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame(cell.text_frame, replacements)

    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for subshape in shape.shapes:
            replace_in_shape(subshape, replacements, question_runs=question_runs)


def fill_slide_placeholders(
    presentation_path: Path,
    output_path: Path,
    slide_replacements: dict[int, dict[str, str]],
    slide_question_runs: dict[int, list[dict[str, bool | str]] | None],
) -> None:
    prs = Presentation(str(presentation_path))

    for slide_number, replacements in slide_replacements.items():
        slide_index = slide_number - 1
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(
                f"В презентации {len(prs.slides)} слайдов, а запрошен слайд {slide_number}."
            )

        slide = prs.slides[slide_index]
        try:
            for shape in slide.shapes:
                replace_in_shape(shape, replacements, question_runs=slide_question_runs.get(slide_number))
        except Exception as exc:
            raise RuntimeError(f"Ошибка при заполнении слайда {slide_number}: {exc}") from exc

    prs.save(str(output_path))


def find_question(questions: list[QuestionItem], number: int) -> QuestionItem | None:
    for question in questions:
        if question.number == number:
            return question
    return None


def get_question_for_number(questions: list[QuestionItem], number: int) -> QuestionItem | None:
    return find_question(questions, number) or (questions[number - 1] if len(questions) >= number else None)


def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Берёт вопросы из Word.docx и подставляет их тематику и текст вопроса "
            "в слайды 6-14,16-24,26-34; «В картинках» 36-41/43-48/50-55; «3х3=12» 57-65/67-75/77-85; «4 Мультимедиа» 89-97/99-107; «Логика и отвага» 109-114/116-121/123-128; «Тематический» 132-137/139-144/146-151; «Блиц» 153-158/160-165"
        )
    )
    parser.add_argument("--word", default="Word.docx", type=Path, help="Путь к Word-файлу")
    parser.add_argument(
        "--template",
        default="Presentation1.pptx",
        type=Path,
        help="Путь к шаблону презентации",
    )
    parser.add_argument(
        "--output",
        default="Presentation1_filled.pptx",
        type=Path,
        help="Путь для сохранения новой презентации",
    )
    args = parser.parse_args()

    questions = parse_questions_from_docx(args.word)
    if not questions:
        raise ValueError("В Word-файле не найдено ни одного корректного блока с Тематикой и Вопросом.")

    max_question_number = 9
    missing_numbers: list[int] = []
    slide_replacements: dict[int, dict[str, str]] = {}
    slide_question_runs: dict[int, list[dict[str, bool | str]] | None] = {}

    for question_number in range(1, max_question_number + 1):
        question = get_question_for_number(questions, question_number)
        if question is None:
            missing_numbers.append(question_number)
            continue

        base_slide_number = question_number + 5
        base_replacements = {
            "тематика": question.theme,
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs
        slide_replacements[base_slide_number + 10] = base_replacements.copy()
        slide_question_runs[base_slide_number + 10] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 20] = answer_slide_replacements
        slide_question_runs[base_slide_number + 20] = question.question_runs

    if missing_numbers:
        missing = ", ".join(str(n) for n in missing_numbers)
        raise ValueError(
            f"В Word-файле не хватает вопросов с номерами: {missing}. "
            "Нужны вопросы №1..№9 для заполнения слайдов 6..14, 16..24 и 26..34."
        )

    round_two_questions = parse_round_without_theme_from_docx(
        args.word,
        round_title="В картинках",
        max_questions=6,
    )
    if len(round_two_questions) < 6:
        raise ValueError(
            "В раунде 'В картинках' найдено меньше 6 вопросов. "
            "Нужны вопросы для слайдов 36..41, 43..48 и 50..55."
        )

    for question_number in range(1, 7):
        question = round_two_questions[question_number - 1]
        base_slide_number = question_number + 35
        base_replacements = {
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs
        slide_replacements[base_slide_number + 7] = base_replacements.copy()
        slide_question_runs[base_slide_number + 7] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 14] = answer_slide_replacements
        slide_question_runs[base_slide_number + 14] = question.question_runs

    round_three_questions = parse_round_with_section_themes_from_docx(
        args.word,
        round_title="3х3=12",
        max_questions=9,
    )
    if len(round_three_questions) < 9:
        raise ValueError(
            "В раунде '3х3=12' найдено меньше 9 вопросов. "
            "Нужны вопросы для слайдов 57..65, 67..75 и 77..85."
        )

    for question_number in range(1, 10):
        question = round_three_questions[question_number - 1]
        base_slide_number = question_number + 56
        base_replacements = {
            "тематика": question.theme,
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs
        slide_replacements[base_slide_number + 10] = base_replacements.copy()
        slide_question_runs[base_slide_number + 10] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 20] = answer_slide_replacements
        slide_question_runs[base_slide_number + 20] = question.question_runs

    round_four_questions = parse_round_without_theme_from_docx(
        args.word,
        round_title="4 Мультимедиа",
        max_questions=9,
    )
    if len(round_four_questions) < 9:
        raise ValueError(
            "В раунде '4 Мультимедиа' найдено меньше 9 вопросов. "
            "Нужны вопросы для слайдов 89..97 и 99..107."
        )

    for question_number in range(1, 10):
        question = round_four_questions[question_number - 1]
        base_slide_number = question_number + 88
        base_replacements = {
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 10] = answer_slide_replacements
        slide_question_runs[base_slide_number + 10] = question.question_runs

    round_five_questions = parse_round_without_theme_from_docx(
        args.word,
        round_title="Логика и отвага",
        max_questions=6,
    )
    if len(round_five_questions) < 6:
        raise ValueError(
            "В раунде 'Логика и отвага' найдено меньше 6 вопросов. "
            "Нужны вопросы для слайдов 109..114, 116..121 и 123..128."
        )

    for question_number in range(1, 7):
        question = round_five_questions[question_number - 1]
        base_slide_number = question_number + 108
        base_replacements = {
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs
        slide_replacements[base_slide_number + 7] = base_replacements.copy()
        slide_question_runs[base_slide_number + 7] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 14] = answer_slide_replacements
        slide_question_runs[base_slide_number + 14] = question.question_runs

    round_six_questions = parse_round_with_constant_theme_from_docx(
        args.word,
        round_title="Тематический",
        theme_label="Тема",
        max_questions=6,
    )
    if len(round_six_questions) < 6:
        raise ValueError(
            "В раунде 'Тематический' найдено меньше 6 вопросов. "
            "Нужны вопросы для слайдов 132..137, 139..144 и 146..151."
        )

    for question_number in range(1, 7):
        question = round_six_questions[question_number - 1]
        base_slide_number = question_number + 131
        base_replacements = {
            "тематика": question.theme,
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs
        slide_replacements[base_slide_number + 7] = base_replacements.copy()
        slide_question_runs[base_slide_number + 7] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 14] = answer_slide_replacements
        slide_question_runs[base_slide_number + 14] = question.question_runs

    round_seven_questions = parse_round_without_theme_from_docx(
        args.word,
        round_title="Блиц",
        max_questions=6,
    )
    if len(round_seven_questions) < 6:
        raise ValueError(
            "В раунде 'Блиц' найдено меньше 6 вопросов. "
            "Нужны вопросы для слайдов 153..158 и 160..165."
        )

    for question_number in range(1, 7):
        question = round_seven_questions[question_number - 1]
        base_slide_number = question_number + 152
        base_replacements = {
            "вопрос": question.question,
        }

        slide_replacements[base_slide_number] = base_replacements.copy()
        slide_question_runs[base_slide_number] = question.question_runs

        answer_slide_replacements = base_replacements.copy()
        answer_slide_replacements["верный ответ"] = question.answer
        slide_replacements[base_slide_number + 7] = answer_slide_replacements
        slide_question_runs[base_slide_number + 7] = question.question_runs

    try:
        fill_slide_placeholders(
            presentation_path=args.template,
            output_path=args.output,
            slide_replacements=slide_replacements,
            slide_question_runs=slide_question_runs,
        )
    except Exception as exc:
        print(f"Ошибка: {exc}")
        raise SystemExit(1) from exc

    print("Готово: заполнение прошло успешно.")


if __name__ == "__main__":
    main()
